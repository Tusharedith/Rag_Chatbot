import pdfplumber, docx, pandas as pd
from pptx import Presentation
from pathlib import Path

def parse_pdf(path: str) -> str:
    with pdfplumber.open(path) as pdf:
        return "\n".join((p.extract_text() or "") for p in pdf.pages)

def parse_docx(path: str) -> str:
    d = docx.Document(path)
    return "\n".join(p.text for p in d.paragraphs if p.text.strip())

def parse_pptx(path: str) -> str:
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)

def parse_table(path: str) -> str:
    if path.lower().endswith(".csv"):
        df = pd.read_csv(path)
    else:
        df = pd.read_excel(path)
    return df.to_csv(index=False)

def parse_txt(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")

def extract_text(path: str) -> str:
    p = path.lower()
    if p.endswith(".pdf"):  return parse_pdf(path)
    if p.endswith(".docx"): return parse_docx(path)
    if p.endswith(".pptx"): return parse_pptx(path)
    if p.endswith(".csv") or p.endswith(".xlsx"): return parse_table(path)
    if p.endswith(".txt"):  return parse_txt(path)
    raise ValueError(f"Unsupported file type: {path}")
